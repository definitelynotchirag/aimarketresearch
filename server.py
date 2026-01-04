import os
import io
import json
import math
import asyncio
from datetime import datetime, timedelta
from typing import Any, Dict, List, Optional, Tuple

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from fastapi import Depends, FastAPI, Header, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse
from pydantic import BaseModel, Field
from aiohttp import ClientSession, ClientTimeout, ClientError
from dotenv import load_dotenv

from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.utils import ImageReader

try:
    from docx import Document
except Exception:
    Document = None


load_dotenv()


class ResearchRequest(BaseModel):
    brand: str = Field(..., description="Brand or company to research")
    country: Optional[str] = Field(None, description="Country or region")
    industry: Optional[str] = Field(None, description="Industry vertical")
    max_results: int = Field(12, ge=1, le=50)
    language: Optional[str] = Field("en")
    use_cache: bool = True


class Competitor(BaseModel):
    name: str
    website: Optional[str] = None
    summary: Optional[str] = None


class SourceItem(BaseModel):
    title: Optional[str] = None
    url: str
    summary: Optional[str] = None


class ResearchResult(BaseModel):
    brand: str
    country: Optional[str]
    industry: Optional[str]
    competitors: List[Competitor]
    trends: List[str] = []
    ad_insights: List[str] = []
    sources: List[SourceItem] = []
    summary: Optional[str] = None


class StrategyRequest(BaseModel):
    brand: str
    industry: Optional[str] = None
    goals: List[str] = ["awareness", "consideration", "leads"]
    budget_usd: float = Field(10000, ge=100)
    country: Optional[str] = None
    research: Optional[ResearchResult] = None
    time_horizon_months: int = Field(3, ge=1, le=12)
    language: Optional[str] = None
    use_cache: bool = True


class StrategyPlan(BaseModel):
    brand: str
    total_budget_usd: float
    allocations: Dict[str, float]
    content_plan: Dict[str, Dict[str, Any]]
    kpis: List[str]
    swot: Dict[str, List[str]]
    timeline: List[Dict[str, Any]]
    media_calendar: List[Dict[str, Any]]


class ReportRequest(BaseModel):
    brand: str
    research: ResearchResult
    strategy: StrategyPlan
    language: Optional[str] = None
    use_cache: bool = True


class WorkflowRequest(BaseModel):
    brand: str
    country: Optional[str] = None
    industry: Optional[str] = None
    goals: List[str] = ["awareness", "consideration", "leads"]
    budget_usd: float = Field(10000, ge=100)
    time_horizon_months: int = Field(3, ge=1, le=12)
    max_results: int = Field(12, ge=1, le=50)


app = FastAPI(title="AI Media Research & Strategy API", version="0.1.0")


def _parse_cors_origins(value: Optional[str]) -> List[str]:
    if not value:
        return ["*"]
    return [v.strip() for v in value.split(",") if v.strip()]


app.add_middleware(
    CORSMiddleware,
    allow_origins=_parse_cors_origins(os.getenv("CORS_ORIGINS")),
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"]
)


async def get_http_session() -> ClientSession:
    session: Optional[ClientSession] = getattr(app.state, "http_session", None)
    if session is None or session.closed:
        # Configure timeout with reasonable values
        timeout = ClientTimeout(
            total=1000,  # 15 minutes total timeout
            connect=1000,  # 30 seconds to establish connection
            sock_read=1000,  # 5 minutes to read data
        )
        session = ClientSession(timeout=timeout)
        app.state.http_session = session
    return session


@app.on_event("shutdown")
async def _shutdown_event() -> None:
    session: Optional[ClientSession] = getattr(app.state, "http_session", None)
    if session and not session.closed:
        await session.close()


def require_api_key(x_api_key: Optional[str] = Header(None)) -> None:
    required_key = os.getenv("API_KEY")
    if required_key and x_api_key != required_key:
        raise HTTPException(status_code=401, detail="Invalid or missing API key")


async def parse_with_groq(content: str, brand: str, country: Optional[str], industry: Optional[str]) -> Optional[Dict[str, Any]]:
    """Use Groq AI to parse Jina's markdown content into structured JSON"""
    api_key = os.getenv("GROQ_API_KEY")
    if not api_key:
        print("GROQ_API_KEY not set, falling back to manual parsing")
        return parse_markdown_research(content, brand, country, industry)
    
    session = await get_http_session()
    
    system_prompt = (
        "You are an expert data parser. Parse the provided market research content and extract structured data. "
        "Return ONLY a valid JSON object with this exact schema:\n"
        "{\n"
        "  \"brand\": string,\n"
        "  \"country\": string|null,\n"
        "  \"industry\": string|null,\n"
        "  \"summary\": string,\n"
        "  \"competitors\": [{ \"name\": string, \"website\": string|null, \"summary\": string }],\n"
        "  \"trends\": [string],\n"
        "  \"ad_insights\": [string],\n"
        "  \"sources\": [{ \"title\": string|null, \"url\": string, \"summary\": string|null }]\n"
        "}\n"
        "Extract competitors from tables and text mentions. Extract market trends and advertising insights. "
        "Include sources from citations. Limit competitors to 12, trends to 8, ad_insights to 6, sources to 12. "
        "Ensure all extracted text is clean and meaningful - no markdown artifacts like '---' or table separators."
    )
    
    user_prompt = (
        f"Parse this market research content for {brand} in {industry or 'the market'}:\n\n"
        f"{content[:8000]}"  # Limit content to avoid token limits
    )
    
    url = os.getenv("GROQ_API_URL", "https://api.groq.com/openai/v1/chat/completions")
    model = os.getenv("GROQ_MODEL", "llama3-70b-8192")
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    
    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        "temperature": 0.1,
        "stream": False,
        "response_format": {"type": "json_object"}
    }
    
    try:
        async with session.post(url, headers=headers, json=payload) as resp:
            if resp.status != 200:
                text = await resp.text()
                print(f"Groq parsing API error: {resp.status} {text}")
                return parse_markdown_research(content, brand, country, industry)
            
            data = await resp.json()
            groq_content = data["choices"][0]["message"]["content"]
            parsed = json.loads(groq_content)
            
            print(f"Groq extracted: {len(parsed.get('competitors', []))} competitors, {len(parsed.get('trends', []))} trends")
            return parsed
            
    except Exception as e:
        print(f"Error using Groq for parsing: {e}")
        return parse_markdown_research(content, brand, country, industry)


def parse_markdown_research(content: str, brand: str, country: Optional[str], industry: Optional[str]) -> Optional[Dict[str, Any]]:
    """Parse Jina's markdown research content into structured JSON format"""
    import re
    
    try:
        # Extract competitors from markdown table
        competitors = []
        
        # Try multiple table patterns
        table_patterns = [
            r'\|\s*Competitor\s*\|\s*Strengths\s*\|\s*Weaknesses\s*\|\s*Market Share.*?\n((?:\|.*?\|.*?\|\n?)+)',
            r'\|\s*Competitor\s*\|\s*Summary\s*\|.*?\n((?:\|.*?\|.*?\|\n?)+)',
            r'Direct competitors include ([^.]+)\.',
            r'competitors include ([^.]+)\.',
        ]
        
        # Try table extraction first
        for pattern in table_patterns[:2]:
            table_match = re.search(pattern, content, re.DOTALL | re.IGNORECASE)
            if table_match:
                rows = table_match.group(1).strip().split('\n')
                for row in rows:
                    # Skip markdown table separators and empty rows
                    if ('|' in row and 
                        not row.strip().startswith('|---') and
                        not row.strip().startswith('| ---') and
                        '---' not in row and
                        len(row.strip()) > 5):
                        
                        parts = [p.strip() for p in row.split('|') if p.strip()]
                        if len(parts) >= 2:
                            name = parts[0].strip()
                            summary = parts[1].strip() if len(parts) > 1 else ""
                            
                            # Comprehensive filtering for valid competitor names
                            if (name and 
                                name.lower() not in ['competitor', 'tesla', '---', '--', '-', 'strengths', 'weaknesses', 'market share'] and 
                                len(name) > 1 and 
                                not name.startswith('-') and
                                not name.endswith('-') and
                                not all(c in '-|=+*#' for c in name) and
                                not re.match(r'^-+$', name) and
                                not re.match(r'^\|.*\|$', name)):
                                
                                competitors.append({
                                    "name": name,
                                    "website": None,
                                    "summary": summary if summary and summary not in ['---', '--', '-'] else f"Competitor in the {industry or 'market'} space"
                                })
                break
        
        # If no table found, extract from text mentions
        if not competitors:
            for pattern in table_patterns[2:]:
                match = re.search(pattern, content, re.IGNORECASE)
                if match:
                    comp_text = match.group(1)
                    # Split by common delimiters
                    comp_names = re.split(r'[,;]\s*(?:and\s+)?', comp_text)
                    for name in comp_names:
                        name = name.strip().replace('and ', '').strip()
                        if name and len(name) > 2:
                            competitors.append({
                                "name": name,
                                "website": None,
                                "summary": f"Competitor in the {industry or 'market'} space"
                            })
                    break
        
        # Extract individual competitor mentions from text
        if not competitors:
            competitor_keywords = [
                r'BYD[^a-zA-Z]',
                r'Ford[^a-zA-Z]',
                r'General Motors',
                r'\bGM\b',
                r'Rivian',
                r'Lucid Motors',
                r'Volkswagen',
                r'BMW[^a-zA-Z]',
                r'NIO[^a-zA-Z]',
                r'Mercedes-Benz',
                r'Toyota[^a-zA-Z]'
            ]
            
            found_competitors = set()
            for keyword in competitor_keywords:
                matches = re.findall(keyword, content, re.IGNORECASE)
                if matches:
                    clean_name = re.sub(r'[^a-zA-Z\s-]', '', keyword.replace('\\b', '').replace('[^a-zA-Z]', ''))
                    found_competitors.add(clean_name.strip())
            
            for name in list(found_competitors)[:12]:
                # Additional filtering for keyword matches
                if (name and 
                    len(name) > 2 and 
                    not name.startswith('-') and
                    name not in ['---', '--', '-', '|'] and
                    not all(c in '-|=' for c in name)):
                    competitors.append({
                        "name": name,
                        "website": None,
                        "summary": f"Major competitor in the {industry or 'automotive'} industry"
                    })
        
        # Extract trends from content
        trends = []
        
        # Look for market trends section
        trend_section_match = re.search(r'### Shifting Market Trends(.*?)(?:###|$)', content, re.DOTALL | re.IGNORECASE)
        if trend_section_match:
            trend_section = trend_section_match.group(1)
            # Extract sentences with key trend indicators
            trend_sentences = re.findall(r'([^.]*(?:market|sales|growth|trend|increasing|declining|competition|technology)[^.]*\.)', trend_section, re.IGNORECASE)
            trends = [t.strip() for t in trend_sentences if len(t.strip()) > 20][:8]
        
        # Fallback: extract key market insights from anywhere in content
        if not trends:
            trend_indicators = [
                r'The US EV market is experiencing ([^.]+\.)',
                r'sales growth ([^.]+\.)',
                r'([^.]*competition[^.]*\.)',
                r'([^.]*market share[^.]*\.)',
                r'([^.]*dominance[^.]*\.)',
                r'([^.]*technology[^.]*\.)'
            ]
            
            for pattern in trend_indicators:
                matches = re.findall(pattern, content, re.IGNORECASE)
                for match in matches[:2]:
                    if len(match.strip()) > 15:
                        trends.append(match.strip())
        
        # Clean and deduplicate trends - filter out markdown artifacts
        cleaned_trends = []
        for t in trends:
            t = t.strip()
            if (len(t) > 15 and 
                not t.startswith('#') and 
                not t.startswith('|') and
                not t.startswith('[^') and
                '---' not in t and
                not all(c in '-|=*#' for c in t[:10])):
                cleaned_trends.append(t)
        trends = list(dict.fromkeys(cleaned_trends))[:8]
        
        # Extract ad insights
        ad_insights = []
        
        # Look for marketing/sales strategy sections
        marketing_sections = [
            r'### Tesla\'s Marketing and Sales Strategies(.*?)(?:###|$)',
            r'marketing strategy(.*?)(?:\[|\n\n)',
            r'direct.*?sales.*?model(.*?)(?:\[|\n\n)',
            r'advertising(.*?)(?:\[|\n\n)'
        ]
        
        for section_pattern in marketing_sections:
            section_match = re.search(section_pattern, content, re.DOTALL | re.IGNORECASE)
            if section_match:
                section_text = section_match.group(1)
                # Extract meaningful insights
                insight_sentences = re.findall(r'([^.]*(?:Tesla|marketing|sales|brand|advertising|direct|consumer|strategy)[^.]*\.)', section_text, re.IGNORECASE)
                ad_insights.extend([s.strip() for s in insight_sentences if len(s.strip()) > 20][:4])
                break
        
        # Fallback: extract marketing-related insights
        if not ad_insights:
            marketing_indicators = [
                r'([^.]*direct.to.consumer[^.]*\.)',
                r'([^.]*brand awareness[^.]*\.)',
                r'([^.]*marketing strategy[^.]*\.)',
                r'([^.]*sales model[^.]*\.)',
                r'([^.]*advertising[^.]*\.)'
            ]
            
            for pattern in marketing_indicators:
                matches = re.findall(pattern, content, re.IGNORECASE)
                for match in matches[:2]:
                    if len(match.strip()) > 15:
                        ad_insights.append(match.strip())
        
        # Clean and deduplicate ad insights - filter out markdown artifacts
        cleaned_insights = []
        for a in ad_insights:
            a = a.strip()
            if (len(a) > 15 and 
                not a.startswith('#') and 
                not a.startswith('|') and
                not a.startswith('[^') and
                '---' not in a and
                not all(c in '-|=*#' for c in a[:10])):
                cleaned_insights.append(a)
        ad_insights = list(dict.fromkeys(cleaned_insights))[:6]
        
        # Extract sources from citations/footnotes
        sources = []
        citation_pattern = r'\[([^\]]+)\]\(([^)]+)\)'
        citations = re.findall(citation_pattern, content)
        for title, url in citations[:12]:
            sources.append({
                "title": title.strip(),
                "url": url.strip(),
                "summary": None
            })
        
        # Extract summary (first few sentences)
        summary_match = re.search(r'^([^.]+\.[^.]+\.)', content.strip())
        summary = summary_match.group(1).strip() if summary_match else f"Research analysis for {brand} in the {industry or 'market'} sector."
        
        # Clean up data
        competitors = competitors[:12]
        trends = list(dict.fromkeys([t for t in trends if len(t) > 10]))[:8]
        ad_insights = list(dict.fromkeys([a for a in ad_insights if len(a) > 10]))[:6]
        
        result = {
            "brand": brand,
            "country": country,
            "industry": industry,
            "summary": summary,
            "competitors": competitors,
            "trends": trends,
            "ad_insights": ad_insights,
            "sources": sources
        }
        
        print(f"Extracted {len(competitors)} competitors, {len(trends)} trends, {len(ad_insights)} insights, {len(sources)} sources")
        print(f"Competitors: {[c['name'] for c in competitors]}")
        print(f"First trend: {trends[0] if trends else 'None'}")
        return result
        
    except Exception as e:
        print(f"Error parsing markdown content: {e}")
        return None


def build_research_query(req: ResearchRequest) -> str:
    parts: List[str] = []
    parts.append(f"Competitors of {req.brand}")
    if req.industry:
        parts.append(f"in industry {req.industry}")
    if req.country:
        parts.append(f"in {req.country}")
    parts.append("pricing, ads, audience, messaging, tone, content strategy, market trends")
    return " ".join(parts)


async def call_jina_deepsearch_structured(req: ResearchRequest) -> Dict[str, Any]:
    api_key = os.getenv("JINA_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="Missing JINA_API_KEY")
    session = await get_http_session()
    if not hasattr(app.state, "jina_lock"):
        app.state.jina_lock = asyncio.Lock()
    if not hasattr(app.state, "jina_cache"):
        app.state.jina_cache = {}
    cache_key = json.dumps({
        "brand": req.brand.strip().lower(),
        "country": (req.country or "").strip().lower(),
        "industry": (req.industry or "").strip().lower(),
        "max_results": req.max_results,
        "language": req.language
    }, sort_keys=True)
    print(cache_key)
    cached = app.state.jina_cache.get(cache_key)
    if req.use_cache and cached:
        cached_time = cached.get("_cached_at")
        if cached_time and (datetime.utcnow() - cached_time).total_seconds() < 60 * 60 * 12:
            return cached["data"]
    system_prompt = (
        "You are an expert market and media research agent. "
        "CRITICAL: Return ONLY a valid JSON object. No explanatory text before or after. No markdown formatting. "
        "Use up-to-date web intelligence to identify direct and indirect competitors, key market trends, and paid media insights. "
        "The JSON schema must be exactly: {\n"
        "  \"brand\": string,\n"
        "  \"country\": string|null,\n"
        "  \"industry\": string|null,\n"
        "  \"summary\": string,\n"
        "  \"competitors\": [ { \"name\": string, \"website\": string|null, \"summary\": string } ],\n"
        "  \"trends\": [ string ],\n"
        "  \"ad_insights\": [ string ],\n"
        "  \"sources\": [ { \"title\": string|null, \"url\": string, \"summary\": string|null } ]\n"
        "}. "
        "Limit competitors to 12 and sources to {max_results}. "
        "Your response must start with {{ and end with }}."
    ).replace("{max_results}", str(req.max_results))
    user_prompt = (
        f"Research scope:\n"
        f"Brand: {req.brand}\n"
        f"Country: {req.country or 'global'}\n"
        f"Industry: {req.industry or 'general'}\n"
        f"Language: {req.language or 'en'}\n"
        f"Task: Identify direct and indirect competitors with brief summaries, extract key market trends, and paid media insights. "
        f"Provide credible sources (title, url, 1-line summary)."
    )
    url = os.getenv("JINA_DEEPSEARCH_URL", "https://deepsearch.jina.ai/v1/chat/completions")
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    payload = {
        "model": "jina-deepsearch-v1",
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        "stream": False,
        "reasoning_effort": "low",
        "temperature": 0.2,
        "response_format": {"type": "json_object"},
        "budget_tokens": 50000,
    }
    print(payload)
    async with app.state.jina_lock:
        async with session.post(url, headers=headers, json=payload) as resp:
            if resp.status != 200:
                text = await resp.text()
                print(f"Jina API Error {resp.status}: {text}")
                raise HTTPException(status_code=502, detail=f"Jina API error: {resp.status} {text}")
            data = await resp.json()

    print(f"Jina API Response: {json.dumps(data, indent=2)}")

    try:
        content = data["choices"][0]["message"]["content"]
    except (KeyError, IndexError, TypeError) as e:
        print(f"Jina response structure error: {e}")
        print(f"Response keys: {list(data.keys()) if isinstance(data, dict) else 'Not a dict'}")
        raise HTTPException(status_code=502, detail=f"Jina API unexpected response format: {str(e)}")

    print(f"Jina content: {content}")

    try:
        parsed = json.loads(content)
    except json.JSONDecodeError as e:
        print(f"JSON parse error: {e}")
        print(f"Content that failed to parse: {content[:500]}...")

        # Use Groq to parse the markdown content into structured JSON
        parsed = await parse_with_groq(content, req.brand, req.country, req.industry)
        if parsed:
            print(f"Successfully parsed markdown content with Groq AI")
        else:
            # Fallback: create a basic response structure
            parsed = {
            "brand": req.brand,
            "country": req.country,
            "industry": req.industry,
            "summary": f"Basic research generated for {req.brand}",
            "competitors": [
                {"name": f"{req.brand} Competitor 1", "website": None, "summary": "Generated competitor"},
                {"name": f"{req.brand} Competitor 2", "website": None, "summary": "Generated competitor"}
            ],
            "trends": ["Market trend 1", "Market trend 2"],
            "ad_insights": ["Ad insight 1", "Ad insight 2"],
            "sources": []
        }
        print(f"Using fallback data: {parsed}")

    app.state.jina_cache[cache_key] = {"data": parsed, "_cached_at": datetime.utcnow()}
    return parsed


async def call_groq_strategy_structured(
    research: ResearchResult,
    goals: List[str],
    budget_usd: float,
    time_horizon_months: int,
    language: Optional[str] = None,
    use_cache: bool = True
) -> Dict[str, Any]:
    api_key = os.getenv("GROQ_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="Missing GROQ_API_KEY")
    session = await get_http_session()
    if not hasattr(app.state, "groq_lock"):
        app.state.groq_lock = asyncio.Lock()
    if not hasattr(app.state, "groq_cache"):
        app.state.groq_cache = {}
    cache_key = json.dumps({
        "brand": research.brand,
        "country": research.country,
        "industry": research.industry,
        "goals": goals,
        "budget_usd": budget_usd,
        "time_horizon_months": time_horizon_months,
        "language": language
    }, sort_keys=True)
    cached = app.state.groq_cache.get(cache_key)
    if use_cache and cached:
        cached_time = cached.get("_cached_at")
        if cached_time and (datetime.utcnow() - cached_time).total_seconds() < 60 * 60 * 12:
            return cached["data"]
    system_prompt = (
        "You are a senior media strategist. Return one valid JSON object only, no markdown. "
        "Use the provided research to produce an actionable media plan. Output schema must be: {\n"
        "  \"allocations\": { \"Google Ads\": number, \"LinkedIn\": number, \"Meta\": number, \"YouTube\": number, \"Twitter/X\": number, \"Others\": number },\n"
        "  \"content_plan\": { \"Video\": { \"cadence_per_week\": number, \"platforms\": [string] }, \"Blog\": { \"cadence_per_week\": number, \"platforms\": [string] }, \"Social\": { \"cadence_per_week\": number, \"platforms\": [string] } },\n"
        "  \"kpis\": [ string ],\n"
        "  \"swot\": { \"strengths\": [string], \"weaknesses\": [string], \"opportunities\": [string], \"threats\": [string] },\n"
        "  \"timeline\": [ { \"phase\": string, \"start\": \"YYYY-MM-DD\", \"end\": \"YYYY-MM-DD\" } ],\n"
        "  \"media_calendar\": [ { \"week\": number, \"channel\": string, \"budget_usd\": number } ]\n"
        "}. "
        "The sum of values in \"allocations\" must equal the total budget {budget}. Use only channels: Google Ads, LinkedIn, Meta, YouTube, Twitter/X, Others. "
        "Make the number of months in the timeline equal to {months}. Respond in the requested language if provided."
    ).replace("{budget}", str(budget_usd)).replace("{months}", str(time_horizon_months))
    research_json = json.dumps(research.dict(), ensure_ascii=False)
    user_prompt = (
        f"Goals: {', '.join(goals)}\n"
        f"Total budget (USD): {budget_usd}\n"
        f"Time horizon (months): {time_horizon_months}\n"
        f"Language: {language or 'en'}\n"
        f"Research JSON: {research_json}\n"
        f"Produce the strategy following the schema with realistic numbers and dates."
    )
    url = os.getenv("GROQ_API_URL", "https://api.groq.com/openai/v1/chat/completions")
    model = os.getenv("GROQ_MODEL", "llama3-70b-8192")
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        "temperature": 0.2,
        "stream": False,
        "response_format": {"type": "json_object"}
    }
    async with app.state.groq_lock:
        async with session.post(url, headers=headers, json=payload) as resp:
            if resp.status != 200:
                text = await resp.text()
                raise HTTPException(status_code=502, detail=f"Groq API error: {resp.status} {text}")
            data = await resp.json()
    try:
        content = data["choices"][0]["message"]["content"]
    except Exception:
        raise HTTPException(status_code=502, detail="Groq API unexpected response format")
    try:
        parsed = json.loads(content)
    except Exception:
        raise HTTPException(status_code=502, detail="Groq API did not return valid JSON content")
    app.state.groq_cache[cache_key] = {"data": parsed, "_cached_at": datetime.utcnow()}
    return parsed


def synthesize_competitors(brand: str, industry: Optional[str], n: int) -> List[Competitor]:
    base = industry or "market"
    names = [
        f"{brand} Labs",
        f"{brand} Media",
        f"{brand} Digital",
        f"{brand} Solutions",
        f"{brand} Edge",
        f"{brand} Analytics",
        f"{brand} Hub",
        f"{brand} Works",
        f"{brand} Pro",
        f"{brand} Nexus"
    ]
    competitors: List[Competitor] = []
    for i in range(min(n, len(names))):
        competitors.append(Competitor(
            name=names[i],
            website=None,
            summary=f"Provider in {base} with focus on growth and performance"
        ))
    return competitors


def parse_jina_results(jina_json: Dict[str, Any], fallback_brand: str, industry: Optional[str], limit: int) -> Tuple[List[Competitor], List[SourceItem], List[str], List[str], str]:
    competitors_data = jina_json.get("competitors") or []
    sources_data = jina_json.get("sources") or []
    trends_data = jina_json.get("trends") or []
    ad_insights_data = jina_json.get("ad_insights") or []
    summary = jina_json.get("summary") or None
    competitors: List[Competitor] = []
    for c in competitors_data[:limit]:
        competitors.append(Competitor(
            name=str(c.get("name") or "").strip() or fallback_brand,
            website=c.get("website"),
            summary=c.get("summary")
        ))
    if not competitors:
        competitors = synthesize_competitors(fallback_brand, industry, limit)
    sources: List[SourceItem] = []
    for s in sources_data[:limit]:
        url = s.get("url")
        if not url:
            continue
        sources.append(SourceItem(title=s.get("title"), url=url, summary=s.get("summary")))
    trends = [str(t) for t in trends_data][:5]
    ad_insights = [str(a) for a in ad_insights_data][:8]
    return competitors, sources, trends, ad_insights, summary or f"Research on {fallback_brand}"


def build_trend_insights(brand: str, industry: Optional[str], sources: List[SourceItem]) -> List[str]:
    base = industry or "the category"
    items = [
        f"Increased competition in {base} with emphasis on multi-channel attribution",
        f"Growing adoption of short-form video for {brand} and peers",
        f"Rise of AI-assisted ad targeting and creative testing",
        f"Greater focus on privacy-safe measurement and first-party data",
        f"Shift in budget towards performance channels and influencer collaborations"
    ]
    return items[:5]


def build_ad_insights(brand: str, industry: Optional[str], sources: List[SourceItem]) -> List[str]:
    base = industry or "the space"
    items = [
        f"Competitors in {base} emphasize Google Search and YouTube for intent and reach",
        f"LinkedIn is effective for B2B lead gen with thought leadership",
        f"Retargeting via Meta improves conversion efficiency",
        f"Content themes center on case studies, ROI proof, and product demos"
    ]
    return items


def allocate_budget(goals: List[str], total_budget: float) -> Dict[str, float]:
    goals_lower = {g.lower() for g in goals}
    if "leads" in goals_lower or "conversions" in goals_lower:
        weights = {"Google Ads": 0.38, "LinkedIn": 0.22, "Meta": 0.18, "YouTube": 0.12, "Twitter/X": 0.05, "Others": 0.05}
    elif "awareness" in goals_lower:
        weights = {"YouTube": 0.3, "Meta": 0.25, "Google Ads": 0.2, "LinkedIn": 0.15, "Twitter/X": 0.05, "Others": 0.05}
    else:
        weights = {"Google Ads": 0.3, "LinkedIn": 0.2, "Meta": 0.2, "YouTube": 0.2, "Others": 0.1}
    allocations = {ch: round(total_budget * w, 2) for ch, w in weights.items()}
    return allocations


def build_content_plan(goals: List[str]) -> Dict[str, Dict[str, Any]]:
    return {
        "Video": {"cadence_per_week": 2, "platforms": ["YouTube", "LinkedIn", "Meta"]},
        "Blog": {"cadence_per_week": 1, "platforms": ["Website", "LinkedIn"]},
        "Social": {"cadence_per_week": 3, "platforms": ["LinkedIn", "Twitter/X", "Meta"]}
    }


def build_kpis(goals: List[str]) -> List[str]:
    goals_lower = {g.lower() for g in goals}
    base = ["CTR", "CPC", "CPM", "Reach", "Engagement Rate"]
    if "leads" in goals_lower:
        base += ["CPL", "SQLs", "Pipeline Value"]
    if "awareness" in goals_lower:
        base += ["Share of Voice", "Brand Search Lift"]
    return list(dict.fromkeys(base))


def build_swot(brand: str, competitors: List[Competitor]) -> Dict[str, List[str]]:
    return {
        "strengths": [f"{brand} product depth", "Agile execution", "Existing customer base"],
        "weaknesses": ["Limited brand awareness", "Small creative library"],
        "opportunities": ["AI-driven optimization", "New channels", "Partnerships"],
        "threats": ["Aggressive competitor pricing", "Platform policy changes"]
    }


def build_timeline(months: int) -> List[Dict[str, Any]]:
    start = datetime.utcnow()
    phases = []
    per_phase = max(1, months // 3)
    labels = ["Discovery & Setup", "Testing & Optimization", "Scale & Expansion"]
    for i, label in enumerate(labels):
        phase_start = start + timedelta(days=30 * per_phase * i)
        phase_end = start + timedelta(days=30 * per_phase * (i + 1))
        phases.append({"phase": label, "start": phase_start.date().isoformat(), "end": phase_end.date().isoformat()})
    return phases


def build_media_calendar(allocations: Dict[str, float]) -> List[Dict[str, Any]]:
    calendar: List[Dict[str, Any]] = []
    weeks = 4
    for i in range(weeks):
        for ch, amt in allocations.items():
            calendar.append({"week": i + 1, "channel": ch, "budget_usd": round(amt / weeks, 2)})
    return calendar


def generate_budget_chart(allocations: Dict[str, float]) -> io.BytesIO:
    labels = list(allocations.keys())
    sizes = list(allocations.values())
    fig, ax = plt.subplots(figsize=(6, 4))
    ax.pie(sizes, labels=labels, autopct='%1.1f%%', startangle=90, counterclock=False)
    ax.axis('equal')
    buf = io.BytesIO()
    plt.tight_layout()
    fig.savefig(buf, format='png', dpi=200)
    plt.close(fig)
    buf.seek(0)
    return buf


def generate_enhanced_budget_chart(allocations: Dict[str, float]) -> io.BytesIO:
    import matplotlib.pyplot as plt
    import matplotlib.patches as patches
    import numpy as np
    
    # Modern color palette
    colors = [
        '#1A237E',  # Deep blue
        '#FF5722',  # Deep orange  
        '#4CAF50',  # Green
        '#FFC107',  # Amber
        '#9C27B0',  # Purple
        '#00BCD4',  # Cyan
        '#E91E63',  # Pink
        '#795548'   # Brown
    ]
    
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(16, 8))
    fig.patch.set_facecolor('#F5F5F5')
    
    channels = list(allocations.keys())
    amounts = list(allocations.values())
    
    # PIE CHART with modern styling
    wedges, texts, autotexts = ax1.pie(
        amounts, 
        labels=channels, 
        colors=colors[:len(channels)],
        autopct='%1.1f%%',
        startangle=90,
        explode=[0.05] * len(channels),  # Slight separation
        shadow=True,
        textprops={'fontsize': 12, 'fontweight': 'bold'}
    )
    
    ax1.set_title('💰 BUDGET DISTRIBUTION', fontsize=18, fontweight='bold', pad=20, color='#1A237E')
    
    # BAR CHART for comparison
    bars = ax2.barh(channels, amounts, color=colors[:len(channels)], alpha=0.8, edgecolor='white', linewidth=2)
    
    # Add value labels on bars
    for i, (bar, amount) in enumerate(zip(bars, amounts)):
        width = bar.get_width()
        ax2.text(width + max(amounts) * 0.01, bar.get_y() + bar.get_height()/2, 
                f'${amount:,.0f}', ha='left', va='center', fontweight='bold', fontsize=11)
    
    ax2.set_title('📊 BUDGET BY CHANNEL', fontsize=18, fontweight='bold', pad=20, color='#1A237E')
    ax2.set_xlabel('Budget ($)', fontsize=14, fontweight='bold')
    ax2.grid(axis='x', alpha=0.3, linestyle='--')
    ax2.set_facecolor('#FAFAFA')
    
    # Style improvements
    for ax in [ax1, ax2]:
        ax.tick_params(labelsize=11)
    
    plt.tight_layout()
    
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=300, bbox_inches='tight', facecolor='#F5F5F5')
    plt.close(fig)
    buf.seek(0)
    return buf


async def call_groq_media_brief_sections(
    research: ResearchResult,
    strategy: StrategyPlan,
    doc_type: str,
    language: Optional[str] = None,
    use_cache: bool = True
) -> Dict[str, Any]:
    api_key = os.getenv("GROQ_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="Missing GROQ_API_KEY")
    session = await get_http_session()
    if not hasattr(app.state, "groq_lock"):
        app.state.groq_lock = asyncio.Lock()
    if not hasattr(app.state, "groq_cache"):
        app.state.groq_cache = {}
    cache_key = json.dumps({
        "doc_type": doc_type,
        "language": language,
        "research": research.dict(),
        "strategy": strategy.dict()
    }, sort_keys=True)
    cached = app.state.groq_cache.get(cache_key)
    if use_cache and cached:
        cached_time = cached.get("_cached_at")
        if cached_time and (datetime.utcnow() - cached_time).total_seconds() < 60 * 60 * 12:
            return cached["data"]
    system_prompt = (
        "You are a senior media strategist and briefing expert. Generate a comprehensive media brief document. Each of the sections should be very big and detailed and rich in content. They should be in depth and comprehensive. They should be very detailed and comprehensive. Return one JSON object only. "
        "Output schema: {\n"
        "  \"client_info\": { \"brand\": string, \"product\": string, \"brief_prepared_by\": string, \"date_of_briefing\": string, \"deadline\": string },\n"
        "  \"marketing_objectives\": string,\n"
        "  \"communication_objectives\": string,\n"
        "  \"tone_manner\": string,\n"
        "  \"media_objectives\": string,\n"
        "  \"brand_positioning\": string,\n"
        "  \"brand_proposition\": string,\n"
        "  \"desired_response\": string,\n"
        "  \"reason_to_believe\": string,\n"
        "  \"target_audience\": { \"demographics\": string, \"psychographics\": string, \"day_in_life\": string },\n"
        "  \"unreached_audience\": string,\n"
        "  \"media_considerations\": string,\n"
        "  \"legal_requirements\": string,\n"
        "  \"creative_considerations\": string,\n"
        "  \"seasonality\": string,\n"
        "  \"geographic_considerations\": string,\n"
        "  \"timing_considerations\": string,\n"
        "  \"budget_provisions\": string,\n"
        "  \"digital_kpis\": string,\n"
        "  \"campaign_duration\": string,\n"
        "  \"creative_assets\": string,\n"
        "  \"website_objectives\": string,\n"
        "  \"conversion_tracking\": string\n"
        "}. Respond in the requested language if provided."
    )
    user_prompt = (
        f"Generate a comprehensive media brief for: {research.brand}\n"
        f"Industry: {research.industry or 'Not specified'}\n"
        f"Target Market: {research.country or 'Global'}\n"
        f"Budget: ${strategy.total_budget_usd:,.0f}\n"
        f"Document type: {doc_type}\n"
        f"Language: {language or 'en'}\n"
        f"Research Data: {json.dumps(research.dict(), ensure_ascii=False)}\n"
        f"Strategy Data: {json.dumps(strategy.dict(), ensure_ascii=False)}\n"
        f"Generate a professional media brief with all required sections populated based on the research and strategy data. "
        f"Use today's date for briefing date and set deadline 30 days from today. Fill in realistic and strategic content for each field."
    )
    url = os.getenv("GROQ_API_URL", "https://api.groq.com/openai/v1/chat/completions")
    model = os.getenv("GROQ_MODEL", "llama3-70b-8192")
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    payload = {
        "model": model,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        "temperature": 0.2,
        "stream": False,
        "response_format": {"type": "json_object"}
    }
    async with app.state.groq_lock:
        async with session.post(url, headers=headers, json=payload) as resp:
            if resp.status != 200:
                text = await resp.text()
                raise HTTPException(status_code=502, detail=f"Groq API error: {resp.status} {text}")
            data = await resp.json()
    try:
        content = data["choices"][0]["message"]["content"]
        parsed = json.loads(content)
    except Exception:
        raise HTTPException(status_code=502, detail="Groq API doc sections invalid format")
    app.state.groq_cache[cache_key] = {"data": parsed, "_cached_at": datetime.utcnow()}
    return parsed


async def build_ppt(brand: str, research: ResearchResult, strategy: StrategyPlan, language: Optional[str] = None, use_cache: bool = True) -> io.BytesIO:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    
    try:
        sections = await call_groq_media_brief_sections(research, strategy, doc_type="pptx", language=language, use_cache=use_cache)
    except HTTPException:
        sections = {}
    
    prs = Presentation()
    
    # Define brand colors - modern gradient scheme
    primary_color = RGBColor(26, 35, 126)     # Deep blue
    accent_color = RGBColor(255, 87, 34)      # Vibrant orange
    success_color = RGBColor(76, 175, 80)     # Success green
    warning_color = RGBColor(255, 193, 7)     # Warning amber
    light_gray = RGBColor(245, 245, 245)      # Light background
    dark_gray = RGBColor(66, 66, 66)          # Dark text
    
    # 🎨 SLIDE 1: MEDIA BRIEF TITLE SLIDE
    title_layout = prs.slide_layouts[6]  # Blank layout for custom design
    slide = prs.slides.add_slide(title_layout)
    
    # Background gradient effect with shapes
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_fill = bg_shape.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = light_gray
    
    # Accent stripe
    accent_stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(2), prs.slide_width
    )
    stripe_fill = accent_stripe.fill
    stripe_fill.solid()
    stripe_fill.fore_color.rgb = primary_color
    
    # Media Brief Title
    title_box = slide.shapes.add_textbox(Inches(2.5), Inches(2), Inches(7), Inches(2))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_para = title_frame.paragraphs[0]
    title_para.text = "MEDIA BRIEF"
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    title_para.alignment = PP_ALIGN.LEFT
    
    # Brand/Product subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(2.5), Inches(3.2), Inches(7), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    subtitle_para = subtitle_frame.paragraphs[0]
    client_info = sections.get("client_info", {}) if isinstance(sections, dict) else {}
    subtitle_para.text = f"Brand: {brand.upper()} | Product: {client_info.get('product', research.industry)}"
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = accent_color
    subtitle_para.alignment = PP_ALIGN.LEFT
    
    # Brief details
    details_box = slide.shapes.add_textbox(Inches(2.5), Inches(4.2), Inches(7), Inches(1.5))
    details_frame = details_box.text_frame
    details_frame.clear()
    details_para = details_frame.paragraphs[0]
    brief_date = client_info.get('date_of_briefing', datetime.utcnow().strftime('%Y-%m-%d'))
    deadline = client_info.get('deadline', (datetime.utcnow() + timedelta(days=30)).strftime('%Y-%m-%d'))
    details_para.text = f"📅 Briefing Date: {brief_date} | ⏰ Deadline: {deadline}"
    details_para.font.size = Pt(14)
    details_para.font.color.rgb = dark_gray
    details_para.alignment = PP_ALIGN.LEFT

    # 🎨 SLIDE 2: MARKETING & COMMUNICATION OBJECTIVES
    layout = prs.slide_layouts[6]  # Blank for custom design
    slide = prs.slides.add_slide(layout)
    
    # Header with gradient background
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header_fill = header_shape.fill
    header_fill.solid()
    header_fill.fore_color.rgb = primary_color
    
    # Title in header
    header_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    header_frame = header_title.text_frame
    header_frame.clear()
    header_para = header_frame.paragraphs[0]
    header_para.text = "1.0 MARKETING & COMMUNICATION OBJECTIVES"
    header_para.font.size = Pt(22)
    header_para.font.bold = True
    header_para.font.color.rgb = RGBColor(255, 255, 255)
    header_para.alignment = PP_ALIGN.CENTER
    
    # Marketing Objectives Section
    obj_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
    obj_frame = obj_box.text_frame
    obj_frame.clear()
    obj_frame.margin_left = Inches(0.2)
    
    marketing_obj_para = obj_frame.paragraphs[0]
    marketing_objectives = sections.get("marketing_objectives", "") if isinstance(sections, dict) else ""
    marketing_obj_para.text = f"1.1 MARKETING OBJECTIVES\n{marketing_objectives or f'Drive brand awareness and market penetration for {brand} in the {research.industry} sector with focus on {research.country} market.'}"
    marketing_obj_para.font.size = Pt(14)
    marketing_obj_para.font.color.rgb = dark_gray
    marketing_obj_para.line_spacing = 1.3
    
    # Communication Objectives Section
    comm_para = obj_frame.add_paragraph()
    comm_objectives = sections.get("communication_objectives", "") if isinstance(sections, dict) else ""
    comm_para.text = f"\n1.2 COMMUNICATION OBJECTIVES\n{comm_objectives or 'Increase brand awareness, improve brand perception, and drive customer consideration through strategic media placement and messaging.'}"
    comm_para.font.size = Pt(14)
    comm_para.font.color.rgb = dark_gray
    comm_para.line_spacing = 1.3
    
    # Tone & Manner Section
    tone_para = obj_frame.add_paragraph()
    tone_manner = sections.get("tone_manner", "") if isinstance(sections, dict) else ""
    tone_para.text = f"\n1.3 TONE & MANNER\n{tone_manner or 'Professional, innovative, and trustworthy communication that resonates with target audience values and aspirations.'}"
    tone_para.font.size = Pt(14)
    tone_para.font.color.rgb = dark_gray
    tone_para.line_spacing = 1.3

    # 🎨 SLIDE 3: BRAND INFORMATION
    slide = prs.slides.add_slide(layout)
    
    # Header
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header_fill = header_shape.fill
    header_fill.solid()
    header_fill.fore_color.rgb = accent_color
    
    header_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    header_frame = header_title.text_frame
    header_frame.clear()
    header_para = header_frame.paragraphs[0]
    header_para.text = "2.0 BRAND INFORMATION"
    header_para.font.size = Pt(24)
    header_para.font.bold = True
    header_para.font.color.rgb = RGBColor(255, 255, 255)
    header_para.alignment = PP_ALIGN.CENTER
    
    # Brand sections
    brand_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    brand_frame = brand_box.text_frame
    brand_frame.clear()
    brand_frame.margin_left = Inches(0.2)
    brand_frame.margin_top = Inches(0.2)
    
    # Brand Positioning
    brand_positioning = sections.get("brand_positioning", "") if isinstance(sections, dict) else ""
    pos_para = brand_frame.paragraphs[0]
    pos_para.text = f"2.1 BRAND POSITIONING\n{brand_positioning or f'{brand} is positioned as a premium, innovative leader in the {research.industry} market, delivering superior value and customer experience.'}"
    pos_para.font.size = Pt(14)
    pos_para.font.color.rgb = dark_gray
    pos_para.line_spacing = 1.3
    
    # Brand Proposition
    brand_proposition = sections.get("brand_proposition", "") if isinstance(sections, dict) else ""
    prop_para = brand_frame.add_paragraph()
    prop_para.text = f"\n2.2 BRAND PROPOSITION\n{brand_proposition or f'Unique value proposition combining cutting-edge technology, superior quality, and exceptional customer service that sets {brand} apart from competitors.'}"
    prop_para.font.size = Pt(14)
    prop_para.font.color.rgb = dark_gray
    prop_para.line_spacing = 1.3
    
    # Desired Response
    desired_response = sections.get("desired_response", "") if isinstance(sections, dict) else ""
    response_para = brand_frame.add_paragraph()
    response_para.text = f"\n2.3 DESIRED RESPONSE\n{desired_response or 'Target customers should perceive our brand as the preferred choice, increase purchase consideration, and actively recommend to others.'}"
    response_para.font.size = Pt(14)
    response_para.font.color.rgb = dark_gray
    response_para.line_spacing = 1.3
    
    # Reason to Believe
    reason_believe = sections.get("reason_to_believe", "") if isinstance(sections, dict) else ""
    reason_para = brand_frame.add_paragraph()
    reason_para.text = f"\n2.4 REASON TO BELIEVE\n{reason_believe or 'Proven track record, customer testimonials, industry awards, and superior product performance validate our brand claims.'}"
    reason_para.font.size = Pt(14)
    reason_para.font.color.rgb = dark_gray
    reason_para.line_spacing = 1.3

    # 🎨 SLIDE 4: TARGET AUDIENCE
    slide = prs.slides.add_slide(layout)
    
    # Header
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header_fill = header_shape.fill
    header_fill.solid()
    header_fill.fore_color.rgb = success_color
    
    header_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    header_frame = header_title.text_frame
    header_frame.clear()
    header_para = header_frame.paragraphs[0]
    header_para.text = "🎯 TARGET AUDIENCE"
    header_para.font.size = Pt(24)
    header_para.font.bold = True
    header_para.font.color.rgb = RGBColor(255, 255, 255)
    header_para.alignment = PP_ALIGN.CENTER
    
    # Target audience sections
    audience_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    audience_frame = audience_box.text_frame
    audience_frame.clear()
    audience_frame.margin_left = Inches(0.2)
    audience_frame.margin_top = Inches(0.2)
    
    # Demographics
    target_audience = sections.get("target_audience", {}) if isinstance(sections, dict) else {}
    demo_para = audience_frame.paragraphs[0]
    demographics = target_audience.get("demographics", "")
    demo_para.text = f"2.2.1 DEMOGRAPHICS\n{demographics or f'Primary audience: {research.country} market, adults 25-54, middle to high income, tech-savvy professionals interested in {research.industry}.'}"
    demo_para.font.size = Pt(14)
    demo_para.font.color.rgb = dark_gray
    demo_para.line_spacing = 1.3
    
    # Psychographics
    psycho_para = audience_frame.add_paragraph()
    psychographics = target_audience.get("psychographics", "")
    psycho_para.text = f"\n2.2.2 PSYCHOGRAPHICS/LIFESTYLE\n{psychographics or 'Values innovation, quality, and sustainability. Early adopters who influence others. Active on social media and research before purchasing.'}"
    psycho_para.font.size = Pt(14)
    psycho_para.font.color.rgb = dark_gray
    psycho_para.line_spacing = 1.3
    
    # Day in the Life
    life_para = audience_frame.add_paragraph()
    day_in_life = target_audience.get("day_in_life", "")
    life_para.text = f"\n2.2.3 DAY IN THE LIFE\n{day_in_life or 'Busy professionals who value efficiency. Start day checking news/social media, work-focused during business hours, evening leisure includes streaming and social browsing.'}"
    life_para.font.size = Pt(14)
    life_para.font.color.rgb = dark_gray
    life_para.line_spacing = 1.3
    
    # Unreached Audience
    unreached_para = audience_frame.add_paragraph()
    unreached_audience = sections.get("unreached_audience", "") if isinstance(sections, dict) else ""
    unreached_para.text = f"\n2.2.4 UNREACHED AUDIENCE\n{unreached_audience or 'Younger demographics (18-25) and emerging markets represent growth opportunities for brand expansion.'}"
    unreached_para.font.size = Pt(14)
    unreached_para.font.color.rgb = dark_gray
    unreached_para.line_spacing = 1.3

    # 🎨 SLIDE 5: EXECUTIONAL CONSIDERATIONS
    slide = prs.slides.add_slide(layout)
    
    # Header
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header_fill = header_shape.fill
    header_fill.solid()
    header_fill.fore_color.rgb = warning_color
    
    header_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    header_frame = header_title.text_frame
    header_frame.clear()
    header_para = header_frame.paragraphs[0]
    header_para.text = "3.0 EXECUTIONAL CONSIDERATIONS"
    header_para.font.size = Pt(22)
    header_para.font.bold = True
    header_para.font.color.rgb = RGBColor(255, 255, 255)
    header_para.alignment = PP_ALIGN.CENTER
    
    # Executional considerations sections
    exec_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    exec_frame = exec_box.text_frame
    exec_frame.clear()
    exec_frame.margin_left = Inches(0.2)
    exec_frame.margin_top = Inches(0.2)
    
    # Media Considerations
    media_considerations = sections.get("media_considerations", "") if isinstance(sections, dict) else ""
    media_para = exec_frame.paragraphs[0]
    media_para.text = f"3.1 MEDIA CONSIDERATIONS\n{media_considerations or f'Focus on digital channels: Google Ads, Social Media, YouTube. Budget allocation: {strategy.allocations}'}"
    media_para.font.size = Pt(12)
    media_para.font.color.rgb = dark_gray
    media_para.line_spacing = 1.2
    
    # Budget Provisions
    budget_provisions = sections.get("budget_provisions", "") if isinstance(sections, dict) else ""
    budget_para = exec_frame.add_paragraph()
    budget_para.text = f"\n3.2 BUDGET PROVISIONS\nTotal Budget: ${strategy.total_budget_usd:,.0f}\n{budget_provisions or 'Budget allocated across multiple channels with flexibility for optimization based on performance metrics.'}"
    budget_para.font.size = Pt(12)
    budget_para.font.color.rgb = dark_gray
    budget_para.line_spacing = 1.2
    
    # Timing Considerations
    timing_considerations = sections.get("timing_considerations", "") if isinstance(sections, dict) else ""
    timing_para = exec_frame.add_paragraph()
    timing_default = f"Campaign launch: {strategy.timeline[0]['start'] if strategy.timeline else 'Q1 2024'}. Timeline: {len(strategy.timeline)} phases over {strategy.timeline[-1]['end'] if strategy.timeline else '6 months'}."
    timing_para.text = f"\n3.3 TIMING CONSIDERATIONS\n{timing_considerations or timing_default}"
    timing_para.font.size = Pt(12)
    timing_para.font.color.rgb = dark_gray
    timing_para.line_spacing = 1.2

    # 🎨 SLIDE 6: DIGITAL CAMPAIGN SPECIFICS
    slide = prs.slides.add_slide(layout)
    
    # Header
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header_fill = header_shape.fill
    header_fill.solid()
    header_fill.fore_color.rgb = RGBColor(156, 39, 176)
    
    header_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    header_frame = header_title.text_frame
    header_frame.clear()
    header_para = header_frame.paragraphs[0]
    header_para.text = "4.0 DIGITAL CAMPAIGN SPECIFICS"
    header_para.font.size = Pt(22)
    header_para.font.bold = True
    header_para.font.color.rgb = RGBColor(255, 255, 255)
    header_para.alignment = PP_ALIGN.CENTER
    
    # Digital campaign sections
    digital_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    digital_frame = digital_box.text_frame
    digital_frame.clear()
    digital_frame.margin_left = Inches(0.2)
    digital_frame.margin_top = Inches(0.2)
    
    # KPIs
    digital_kpis = sections.get("digital_kpis", "") if isinstance(sections, dict) else ""
    kpi_para = digital_frame.paragraphs[0]
    kpi_para.text = f"4.1 KEY PERFORMANCE INDICATORS\n{digital_kpis or f'Primary KPIs: {', '.join(strategy.kpis[:5])}'}"
    kpi_para.font.size = Pt(14)
    kpi_para.font.color.rgb = dark_gray
    kpi_para.line_spacing = 1.3
    
    # Campaign Duration
    campaign_duration = sections.get("campaign_duration", "") if isinstance(sections, dict) else ""
    duration_para = digital_frame.add_paragraph()
    duration_para.text = f"\n4.2 CAMPAIGN DURATION\n{campaign_duration or f'Campaign runs for {len(strategy.timeline)} phases, approximately {len(strategy.media_calendar) // 4} months with continuous optimization.'}"
    duration_para.font.size = Pt(14)
    duration_para.font.color.rgb = dark_gray
    duration_para.line_spacing = 1.3
    
    # Website & Tracking
    website_objectives = sections.get("website_objectives", "") if isinstance(sections, dict) else ""
    conversion_tracking = sections.get("conversion_tracking", "") if isinstance(sections, dict) else ""
    web_para = digital_frame.add_paragraph()
    web_para.text = f"\n4.3 WEBSITE & CONVERSION TRACKING\n{website_objectives or 'Drive traffic to optimized landing pages focused on conversion.'}\n{conversion_tracking or 'Implement Google Analytics, conversion pixels, and attribution tracking across all channels.'}"
    web_para.font.size = Pt(14)
    web_para.font.color.rgb = dark_gray
    web_para.line_spacing = 1.3

    # 🎨 SLIDE 7: MEDIA BRIEF SUMMARY
    slide = prs.slides.add_slide(layout)
    
    # Gradient background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_fill = bg_shape.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = primary_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_para = title_frame.paragraphs[0]
    title_para.text = "MEDIA BRIEF COMPLETE"
    title_para.font.size = Pt(42)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Summary
    summary_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(2))
    summary_frame = summary_box.text_frame
    summary_frame.clear()
    summary_para = summary_frame.paragraphs[0]
    summary_para.text = f"Ready to execute strategic media campaign for {brand}\nBudget: ${strategy.total_budget_usd:,.0f} | Timeline: {len(strategy.timeline)} phases\nTarget: {research.country} {research.industry} market"
    summary_para.font.size = Pt(18)
    summary_para.font.color.rgb = accent_color
    summary_para.alignment = PP_ALIGN.CENTER
    summary_para.line_spacing = 1.4

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


async def build_pdf(brand: str, research: ResearchResult, strategy: StrategyPlan, language: Optional[str] = None, use_cache: bool = True) -> io.BytesIO:
    try:
        sections = await call_groq_media_brief_sections(research, strategy, doc_type="pdf", language=language, use_cache=use_cache)
    except HTTPException:
        sections = {}
    
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4)
    styles = getSampleStyleSheet()
    story: List[Any] = []
    
    # Title and Client Info
    client_info = sections.get("client_info", {}) if isinstance(sections, dict) else {}
    story.append(Paragraph("MEDIA BRIEF", styles["Title"]))
    story.append(Spacer(1, 12))
    story.append(Paragraph(f"Client: {brand}", styles["Heading2"]))
    story.append(Paragraph(f"Brand/Product: {client_info.get('product', research.industry)}", styles["BodyText"]))
    story.append(Paragraph(f"Brief Prepared By: AI Media Research", styles["BodyText"]))
    story.append(Paragraph(f"Date of Briefing: {client_info.get('date_of_briefing', datetime.now().strftime('%Y-%m-%d'))}", styles["BodyText"]))
    story.append(Paragraph(f"Deadline: {client_info.get('deadline', (datetime.now() + timedelta(days=30)).strftime('%Y-%m-%d'))}", styles["BodyText"]))
    story.append(Spacer(1, 12))
    
    # Marketing Objectives
    story.append(Paragraph("1.0 MARKETING OBJECTIVES", styles["Heading2"]))
    marketing_objectives = sections.get("marketing_objectives", "") if isinstance(sections, dict) else ""
    story.append(Paragraph(marketing_objectives or f"Drive brand awareness and market penetration for {brand} in the {research.industry} sector with focus on {research.country} market.", styles["BodyText"]))
    story.append(Spacer(1, 12))
    
    # Communication Objectives  
    story.append(Paragraph("1.1 COMMUNICATION OBJECTIVES", styles["Heading2"]))
    comm_objectives = sections.get("communication_objectives", "") if isinstance(sections, dict) else ""
    story.append(Paragraph(comm_objectives or "Increase brand awareness, improve brand perception, and drive customer consideration through strategic media placement and messaging.", styles["BodyText"]))
    story.append(Spacer(1, 12))
    
    # Brand Information
    story.append(Paragraph("2.0 BRAND INFORMATION", styles["Heading2"]))
    brand_positioning = sections.get("brand_positioning", "") if isinstance(sections, dict) else ""
    story.append(Paragraph("2.1 Brand Positioning:", styles["Heading3"]))
    story.append(Paragraph(brand_positioning or f"{brand} is positioned as a premium, innovative leader in the {research.industry} market, delivering superior value and customer experience.", styles["BodyText"]))
    story.append(Spacer(1, 6))
    
    brand_proposition = sections.get("brand_proposition", "") if isinstance(sections, dict) else ""
    story.append(Paragraph("2.2 Brand Proposition:", styles["Heading3"]))
    story.append(Paragraph(brand_proposition or f"Unique value proposition combining cutting-edge technology, superior quality, and exceptional customer service that sets {brand} apart from competitors.", styles["BodyText"]))
    story.append(Spacer(1, 12))
    
    # Target Audience
    story.append(Paragraph("2.2 TARGET AUDIENCE", styles["Heading2"]))
    target_audience = sections.get("target_audience", {}) if isinstance(sections, dict) else {}
    demographics = target_audience.get("demographics", "")
    story.append(Paragraph("Demographics:", styles["Heading3"]))
    story.append(Paragraph(demographics or f"Primary audience: {research.country} market, adults 25-54, middle to high income, tech-savvy professionals interested in {research.industry}.", styles["BodyText"]))
    story.append(Spacer(1, 12))
    
    # Budget Allocation
    story.append(Paragraph("3.0 BUDGET PROVISIONS", styles["Heading2"]))
    story.append(Paragraph(f"Total Budget: ${strategy.total_budget_usd:,.0f}", styles["BodyText"]))
    chart_buf = generate_budget_chart(strategy.allocations)
    image_reader = ImageReader(chart_buf)
    story.append(RLImage(image_reader, width=400, height=260))
    story.append(Spacer(1, 12))
    
    # Digital Campaign Specifics
    story.append(Paragraph("4.0 DIGITAL CAMPAIGN SPECIFICS", styles["Heading2"]))
    digital_kpis = sections.get("digital_kpis", "") if isinstance(sections, dict) else ""
    story.append(Paragraph("KPIs:", styles["Heading3"]))
    story.append(Paragraph(digital_kpis or f"Primary KPIs: {', '.join(strategy.kpis[:5])}", styles["BodyText"]))
    story.append(Spacer(1, 6))
    
    campaign_duration = sections.get("campaign_duration", "") if isinstance(sections, dict) else ""
    story.append(Paragraph("Campaign Duration:", styles["Heading3"]))
    story.append(Paragraph(campaign_duration or f"Campaign runs for {len(strategy.timeline)} phases, approximately {len(strategy.media_calendar) // 4} months with continuous optimization.", styles["BodyText"]))
    story.append(Spacer(1, 12))
    
    doc.build(story)
    buf.seek(0)
    return buf


async def build_docx(brand: str, research: ResearchResult, strategy: StrategyPlan, language: Optional[str] = None, use_cache: bool = True) -> io.BytesIO:
    if Document is None:
        raise HTTPException(status_code=501, detail="python-docx not installed")
    try:
        sections = await call_groq_media_brief_sections(research, strategy, doc_type="docx", language=language, use_cache=use_cache)
    except HTTPException:
        sections = {}
    
    doc = Document()
    doc.add_heading("MEDIA BRIEF", 0)
    
    # Client Information
    client_info = sections.get("client_info", {}) if isinstance(sections, dict) else {}
    doc.add_heading("Client Information", level=1)
    doc.add_paragraph(f"Client: {brand}")
    doc.add_paragraph(f"Brand/Product: {client_info.get('product', research.industry)}")
    doc.add_paragraph(f"Brief Prepared By: AI Media Research")
    doc.add_paragraph(f"Date of Briefing: {client_info.get('date_of_briefing', datetime.now().strftime('%Y-%m-%d'))}")
    doc.add_paragraph(f"Deadline: {client_info.get('deadline', (datetime.now() + timedelta(days=30)).strftime('%Y-%m-%d'))}")
    
    # Marketing Objectives
    doc.add_heading("1.0 Marketing Objectives", level=1)
    marketing_objectives = sections.get("marketing_objectives", "") if isinstance(sections, dict) else ""
    doc.add_paragraph(marketing_objectives or f"Drive brand awareness and market penetration for {brand} in the {research.industry} sector with focus on {research.country} market.")
    
    # Communication Objectives
    doc.add_heading("1.1 Communication Objectives", level=2)
    comm_objectives = sections.get("communication_objectives", "") if isinstance(sections, dict) else ""
    doc.add_paragraph(comm_objectives or "Increase brand awareness, improve brand perception, and drive customer consideration through strategic media placement and messaging.")
    
    # Tone & Manner
    doc.add_heading("1.2 Tone & Manner of Communication", level=2)
    tone_manner = sections.get("tone_manner", "") if isinstance(sections, dict) else ""
    doc.add_paragraph(tone_manner or "Professional, innovative, and trustworthy communication that resonates with target audience values and aspirations.")
    
    # Brand Information
    doc.add_heading("2.0 Brand Information", level=1)
    brand_positioning = sections.get("brand_positioning", "") if isinstance(sections, dict) else ""
    doc.add_heading("2.1 Brand Positioning", level=2)
    doc.add_paragraph(brand_positioning or f"{brand} is positioned as a premium, innovative leader in the {research.industry} market, delivering superior value and customer experience.")
    
    brand_proposition = sections.get("brand_proposition", "") if isinstance(sections, dict) else ""
    doc.add_heading("2.2 Brand Proposition / Benefit", level=2)
    doc.add_paragraph(brand_proposition or f"Unique value proposition combining cutting-edge technology, superior quality, and exceptional customer service that sets {brand} apart from competitors.")
    
    # Target Audience
    doc.add_heading("2.2 Target Audience", level=1)
    target_audience = sections.get("target_audience", {}) if isinstance(sections, dict) else {}
    
    demographics = target_audience.get("demographics", "")
    doc.add_heading("Demographics", level=2)
    doc.add_paragraph(demographics or f"Primary audience: {research.country} market, adults 25-54, middle to high income, tech-savvy professionals interested in {research.industry}.")
    
    psychographics = target_audience.get("psychographics", "")
    doc.add_heading("Psychographics / Lifestyle", level=2)
    doc.add_paragraph(psychographics or "Values innovation, quality, and sustainability. Early adopters who influence others. Active on social media and research before purchasing.")
    
    # Executional Considerations
    doc.add_heading("3.0 Executional Considerations", level=1)
    media_considerations = sections.get("media_considerations", "") if isinstance(sections, dict) else ""
    doc.add_heading("Media Considerations", level=2)
    doc.add_paragraph(media_considerations or f"Focus on digital channels: Google Ads, Social Media, YouTube. Budget allocation: {strategy.allocations}")
    
    doc.add_heading("Budget Provisions", level=2)
    budget_provisions = sections.get("budget_provisions", "") if isinstance(sections, dict) else ""
    doc.add_paragraph(f"Total Budget: ${strategy.total_budget_usd:,.0f}")
    doc.add_paragraph(budget_provisions or "Budget allocated across multiple channels with flexibility for optimization based on performance metrics.")
    
    # Budget Chart
    chart_buf = generate_budget_chart(strategy.allocations)
    chart_buf.name = "budget.png"
    doc.add_picture(chart_buf, width=None)
    
    # Digital Campaign Specifics
    doc.add_heading("4.0 Digital/Online Campaign Only", level=1)
    digital_kpis = sections.get("digital_kpis", "") if isinstance(sections, dict) else ""
    doc.add_heading("KPI", level=2)
    doc.add_paragraph(digital_kpis or f"Primary KPIs: {', '.join(strategy.kpis[:5])}")
    
    campaign_duration = sections.get("campaign_duration", "") if isinstance(sections, dict) else ""
    doc.add_heading("Campaign Duration", level=2)
    doc.add_paragraph(campaign_duration or f"Campaign runs for {len(strategy.timeline)} phases, approximately {len(strategy.media_calendar) // 4} months with continuous optimization.")
    
    website_objectives = sections.get("website_objectives", "") if isinstance(sections, dict) else ""
    doc.add_heading("Website & Objectives", level=2)
    doc.add_paragraph(website_objectives or "Drive traffic to optimized landing pages focused on conversion.")
    
    conversion_tracking = sections.get("conversion_tracking", "") if isinstance(sections, dict) else ""
    doc.add_heading("Conversion Tracking", level=2)
    doc.add_paragraph(conversion_tracking or "Implement Google Analytics, conversion pixels, and attribution tracking across all channels.")
    
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


@app.get("/health")
async def health() -> Dict[str, Any]:
    return {"status": "ok", "time": datetime.utcnow().isoformat()}


@app.post("/test-parsing")
async def test_parsing(test_data: Dict[str, Any]) -> Dict[str, Any]:
    """Test the markdown parsing function with example Jina response"""
    try:
        # Extract the content from the test data
        content = test_data.get("content", "")
        brand = test_data.get("brand", "TestBrand")
        country = test_data.get("country", "TestCountry")
        industry = test_data.get("industry", "TestIndustry")
        
        if not content:
            return {"error": "No content provided. Send JSON with 'content' field containing the Jina response text."}
        
        # Test the parsing function with Groq
        result = await parse_with_groq(content, brand, country, industry)
        
        if result:
            return {
                "status": "success",
                "message": "Parsing completed successfully",
                "extracted_data": result,
                "stats": {
                    "competitors_found": len(result.get("competitors", [])),
                    "trends_found": len(result.get("trends", [])),
                    "ad_insights_found": len(result.get("ad_insights", [])),
                    "sources_found": len(result.get("sources", []))
                }
            }
        else:
            return {
                "status": "failed",
                "message": "Parsing failed - check server logs for details",
                "extracted_data": None
            }
            
    except Exception as e:
        return {
            "status": "error", 
            "message": f"Error during parsing test: {str(e)}",
            "extracted_data": None
        }


@app.get("/test-jina")
async def test_jina() -> Dict[str, Any]:
    api_key = os.getenv("JINA_API_KEY")
    if not api_key:
        return {"status": "error", "message": "JINA_API_KEY not set"}
    
    if api_key == "test":
        return {"status": "warning", "message": "JINA_API_KEY is set to 'test' - please set a real API key"}
    
    session = await get_http_session()
    url = os.getenv("JINA_DEEPSEARCH_URL", "https://deepsearch.jina.ai/v1/chat/completions")
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    
    # Simple test payload
    payload = {
        "model": "jina-deepsearch-v1",
        "messages": [{"role": "user", "content": "Hello, can you respond with just 'OK'?"}],
        "stream": False,
        "reasoning_effort": "low",
        "temperature": 0.1
    }
    
    try:
        async with session.post(url, headers=headers, json=payload) as resp:
            status = resp.status
            if status == 401:
                return {"status": "error", "message": "Invalid Jina API key"}
            elif status != 200:
                text = await resp.text()
                return {"status": "error", "message": f"Jina API error: {status} {text}"}
            
            data = await resp.json()
            return {"status": "success", "message": "Jina API connection successful", "response_keys": list(data.keys())}
    
    except Exception as e:
        return {"status": "error", "message": f"Connection error: {str(e)}"}


@app.post("/api/research", dependencies=[Depends(require_api_key)])
async def research(req: ResearchRequest) -> ResearchResult:
    jina_json = await call_jina_deepsearch_structured(req)
    competitors, sources, trends, ad_insights, summary = parse_jina_results(jina_json, req.brand, req.industry, req.max_results)
    return ResearchResult(
        brand=req.brand,
        country=req.country,
        industry=req.industry,
        competitors=competitors,
        trends=trends,
        ad_insights=ad_insights,
        sources=sources,
        summary=summary
    )


@app.post("/api/strategy", dependencies=[Depends(require_api_key)])
async def strategy(req: StrategyRequest) -> StrategyPlan:
    research_data: Optional[ResearchResult] = req.research
    if research_data is None:
        raise HTTPException(status_code=400, detail="'research' is required. Call /api/research or /api/workflow first.")
    try:
        groq_json = await call_groq_strategy_structured(
            research=research_data,
            goals=req.goals,
            budget_usd=req.budget_usd,
            time_horizon_months=req.time_horizon_months,
            language=req.language or "en",
            use_cache=req.use_cache
        )
    except HTTPException:
        groq_json = {}
    allocations = groq_json.get("allocations") or allocate_budget(req.goals, req.budget_usd)
    total = sum(float(v) for v in allocations.values()) if isinstance(allocations, dict) else 0.0
    if isinstance(allocations, dict) and total > 0 and abs(total - req.budget_usd) > 0.01:
        factor = req.budget_usd / total
        allocations = {k: round(float(v) * factor, 2) for k, v in allocations.items()}
    content_plan = groq_json.get("content_plan") or build_content_plan(req.goals)
    kpis = groq_json.get("kpis") or build_kpis(req.goals)
    swot = groq_json.get("swot") or build_swot(req.brand, research_data.competitors)
    timeline = groq_json.get("timeline") or build_timeline(req.time_horizon_months)
    calendar = groq_json.get("media_calendar") or build_media_calendar(allocations)
    return StrategyPlan(
        brand=req.brand,
        total_budget_usd=req.budget_usd,
        allocations=allocations,
        content_plan=content_plan,
        kpis=kpis,
        swot=swot,
        timeline=timeline,
        media_calendar=calendar
    )


@app.post("/api/generate/pptx", dependencies=[Depends(require_api_key)])
async def generate_pptx(req: ReportRequest) -> StreamingResponse:
    buf = await build_ppt(req.brand, req.research, req.strategy, language=req.language, use_cache=req.use_cache)
    filename = f"{req.brand.replace(' ', '_').lower()}_strategy_{datetime.utcnow().strftime('%Y%m%d')}.pptx"
    return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={"Content-Disposition": f"attachment; filename={filename}"})


@app.post("/api/generate/pdf", dependencies=[Depends(require_api_key)])
async def generate_pdf(req: ReportRequest) -> StreamingResponse:
    buf = await build_pdf(req.brand, req.research, req.strategy, language=req.language, use_cache=req.use_cache)
    filename = f"{req.brand.replace(' ', '_').lower()}_report_{datetime.utcnow().strftime('%Y%m%d')}.pdf"
    return StreamingResponse(buf, media_type="application/pdf", headers={"Content-Disposition": f"attachment; filename={filename}"})


@app.post("/api/generate/docx", dependencies=[Depends(require_api_key)])
async def generate_docx(req: ReportRequest) -> StreamingResponse:
    buf = await build_docx(req.brand, req.research, req.strategy, language=req.language, use_cache=req.use_cache)
    filename = f"{req.brand.replace(' ', '_').lower()}_report_{datetime.utcnow().strftime('%Y%m%d')}.docx"
    return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", headers={"Content-Disposition": f"attachment; filename={filename}"})


@app.post("/api/workflow", dependencies=[Depends(require_api_key)])
async def workflow(req: WorkflowRequest) -> Dict[str, Any]:
    rreq = ResearchRequest(brand=req.brand, country=req.country, industry=req.industry, max_results=req.max_results, use_cache=True)
    research_result = await research(rreq)
    sreq = StrategyRequest(brand=req.brand, industry=req.industry, goals=req.goals, budget_usd=req.budget_usd, country=req.country, research=research_result, time_horizon_months=req.time_horizon_months)
    strategy_plan = await strategy(sreq)
    return {"research": research_result.dict(), "strategy": strategy_plan.dict()}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("server:app", host="0.0.0.0", port=int(os.getenv("PORT", "8000")), reload=True)
